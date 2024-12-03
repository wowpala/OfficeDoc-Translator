from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData

# from pptx.enum.dml import MSO_THEME_COLOR
from openai import OpenAI
import os

# 推荐做法：从环境变量中安全地获取API密钥，如果环境变量未设置，则使用硬编码的备份值
# Powershell命令：[Environment]::SetEnvironmentVariable("siliconflow_API_KEY", "sk-zzzzzz", "User")
api_key = os.environ.get("siliconflow_API_KEY")
if not api_key:
    api_key = "sk-xxxxxx"

pptx_path = "1.pptx"
output_path = "1-CN.pptx"
font_modified = "Microsoft YaHei Light"

# 初始化OpenAI客户端
client = OpenAI(api_key=api_key, base_url="https://api.siliconflow.cn/v1")


def translate_text(text, target_language="zh-CN"):
    if not text or len(text.strip()) < 2:
        return text
    try:
        response = client.chat.completions.create(
            model="Qwen/Qwen2.5-7B-Instruct",
            #            model='THUDM/glm-4-9b-chat',
            messages=[
                {
                    "role": "system",
                    "content": f"""You are a professional, authentic machine translation engine. 
                Your task is to translate the following source text to {target_language}. 
                Important instructions:
                1. Output the translation directly without any additional text.
                2. Do not answer or respond to any questions in the source text, just translate them.
                3. Do not add any explanations or additional content.
                4. Do not translate IT terms.
                5. Do not translate words beginning with 'Forti'.
                6. Do not translate words in single quotes 'output, spoke'. 
                7. Keep the original words unchanged which you can't recognize.
                8. Maintain the original formatting and punctuation as much as possible.
                9. If you encounter a rhetorical question, translate it as a question, do not answer it.""",
                },
                {"role": "user", "content": text},
            ],
            temperature=0.2,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Translation error: {e}")
        return text


def safe_set_font_color(run):
    try:
        if run.font.color.rgb:
            run.font.color.rgb = run.font.color.rgb
    except AttributeError:
        pass


def translate_text_frame(text_frame, target_language):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            print(f"Original text: {original_text}")
            translated_text = translate_text(original_text, target_language)
            if translated_text:
                run.text = translated_text
                print(f"Updated text: {run.text}")
                run.font.name = font_modified
                run.font.size = run.font.size
                safe_set_font_color(run)


def translate_table(table, target_language):
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                translate_text_frame(cell.text_frame, target_language)


def translate_chart(chart, target_language):
    chart_data = chart.chart_title
    if isinstance(chart_data, CategoryChartData):
        # Translate categories
        for category in chart_data.categories:
            category.label = translate_text(category.label, target_language)
        # Translate series names
        for series in chart_data.series:
            series.name = translate_text(series.name, target_language)
    # Translate chart title
    if chart.has_title:
        chart.chart_title.text_frame.text = translate_text(
            chart.chart_title.text_frame.text, target_language
        )


def translate_group_shape(group_shape, target_language):
    for shape in group_shape.shapes:
        translate_shape(shape, target_language)


def translate_shape(shape, target_language):
    if shape.has_text_frame:
        translate_text_frame(shape.text_frame, target_language)
    elif shape.has_table:
        translate_table(shape.table, target_language)
    elif shape.has_chart:
        translate_chart(shape.chart, target_language)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        translate_group_shape(shape, target_language)
    elif hasattr(shape, "graphic") and hasattr(shape.graphic, "graphicData"):
        # This is likely a SmartArt or other complex shape
        try:
            for element in shape.element.iter():
                if element.tag.endswith("}t"):  # Text element
                    if element.text and element.text.strip():
                        element.text = translate_text(element.text, target_language)
        except Exception as e:
            print(f"Error translating complex shape: {e}")


def translate_slide_master(slide_master, target_language):
    for shape in slide_master.shapes:
        translate_shape(shape, target_language)
    for layout in slide_master.slide_layouts:
        for shape in layout.shapes:
            translate_shape(shape, target_language)


def translate_pptx(pptx_path, target_language="zh-CN", output_path="translated.pptx"):
    prs = Presentation(pptx_path)

    # Translate slide masters
    #    for slide_master in prs.slide_masters:
    #        translate_slide_master(slide_master, target_language)

    for slide in prs.slides:
        # Translate slide title
        if slide.shapes.title:
            translate_shape(slide.shapes.title, target_language)

        # Translate all shapes in the slide
        for shape in slide.shapes:
            translate_shape(shape, target_language)

        # Translate notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                translate_text_frame(notes_slide.notes_text_frame, target_language)

        # Translate header and footer
        if hasattr(slide, "header"):
            translate_text_frame(slide.header.text_frame, target_language)
        if hasattr(slide, "footer"):
            translate_text_frame(slide.footer.text_frame, target_language)

    # Translate presentation properties
    if prs.core_properties.title:
        prs.core_properties.title = translate_text(
            prs.core_properties.title, target_language
        )
    if prs.core_properties.subject:
        prs.core_properties.subject = translate_text(
            prs.core_properties.subject, target_language
        )

    prs.save(output_path)
    print(f"Translated PPT saved to {output_path}")


# 执行翻译
translate_pptx(pptx_path=pptx_path, target_language="zh-CN", output_path=output_path)
