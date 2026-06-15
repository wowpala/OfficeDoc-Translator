import argparse
import os
import tempfile
import unittest

import requests
from pptx import Presentation

import OfficeDoc_Translator as translator


class FakeResponse:
    def __init__(self, status_code=200, json_data=None, text="", headers=None):
        self.status_code = status_code
        self._json_data = json_data or {}
        self.text = text
        self.headers = headers or {}

    def raise_for_status(self):
        if self.status_code >= 400:
            raise requests.exceptions.HTTPError(f"HTTP {self.status_code}")

    def json(self):
        return self._json_data


class SequenceSession:
    def __init__(self, responses):
        self.responses = list(responses)

    def post(self, *args, **kwargs):
        response = self.responses.pop(0)
        if isinstance(response, Exception):
            raise response
        return response


class FakeClient:
    def __init__(self):
        self.calls = []

    def translate(self, text, target_language):
        self.calls.append((text, target_language))
        return f"ZH:{text}"


class FailingClient:
    def __init__(self, message):
        self.message = message

    def translate(self, text, target_language):
        raise translator.TranslationRequestError(self.message)


class OfficeDocTranslatorTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.original_cache_dir = translator.CACHE_DIR
        translator.CACHE_DIR = self.temp_dir.name

    def tearDown(self):
        translator.CACHE_DIR = self.original_cache_dir
        self.temp_dir.cleanup()

    def make_config(self, **overrides):
        values = {
            "api_key": "test-key",
            "model_name": "test-model",
            "endpoint": "https://api.groq.com/openai/v1",
            "temperature": 1.0,
            "enable_thinking": False,
            "request_timeout_seconds": 5.0,
            "max_retries": 2,
            "initial_backoff_seconds": 1.0,
            "max_backoff_seconds": 5.0,
            "min_request_interval_seconds": 0.0,
            "consecutive_failure_limit": 3,
            "save_progress_every_n": 2,
            "resume_enabled": True,
        }
        values.update(overrides)
        return translator.Config(**values)

    def make_runtime(self, input_path, **arg_overrides):
        args = argparse.Namespace(
            no_cache=arg_overrides.get("no_cache", True),
            resume=arg_overrides.get("resume", False),
            fail_fast=arg_overrides.get("fail_fast", True),
        )
        config = self.make_config(save_progress_every_n=1)
        return translator.TranslationRuntime(
            args=args,
            config=config,
            input_file=input_path,
            target_language="zh-CN",
        )

    def test_llmclient_retries_429_and_honors_retry_after(self):
        sleeps = []
        session = SequenceSession(
            [
                FakeResponse(status_code=429, text="too many", headers={"Retry-After": "2"}),
                FakeResponse(
                    status_code=200,
                    json_data={"choices": [{"message": {"content": "数据保护"}}]},
                ),
            ]
        )
        client = translator.LLMClient(
            config=self.make_config(max_retries=1),
            fail_fast=True,
            session=session,
            sleep_fn=sleeps.append,
            monotonic_fn=lambda: 0.0,
            random_fn=lambda: 0.0,
        )

        result = client.translate("Data Protection", "zh-CN")

        self.assertEqual("数据保护", result)
        self.assertEqual(2, client.api_request_count)
        self.assertEqual(1, client.retry_count)
        self.assertEqual([2.0], sleeps)

    def test_llmclient_retries_ssl_then_aborts(self):
        sleeps = []
        session = SequenceSession(
            [
                requests.exceptions.SSLError("ssl eof"),
                requests.exceptions.SSLError("ssl eof"),
            ]
        )
        client = translator.LLMClient(
            config=self.make_config(max_retries=1),
            fail_fast=True,
            session=session,
            sleep_fn=sleeps.append,
            monotonic_fn=lambda: 0.0,
            random_fn=lambda: 0.0,
        )

        with self.assertRaises(translator.TranslationAbort):
            client.translate("Inline CASB", "zh-CN")

        self.assertEqual(2, client.api_request_count)
        self.assertEqual(1, client.retry_count)
        self.assertEqual([1.0], sleeps)

    def test_runtime_reuses_same_run_cache_when_persistence_disabled(self):
        input_path = os.path.join(self.temp_dir.name, "sample.docx")
        with open(input_path, "wb") as handle:
            handle.write(b"doc")

        runtime = self.make_runtime(input_path, no_cache=True, resume=False, fail_fast=True)
        fake_client = FakeClient()
        runtime.client = fake_client

        first = runtime.translate_text("Data Protection")
        second = runtime.translate_text("Data Protection")

        self.assertEqual("ZH:Data Protection", first)
        self.assertEqual("ZH:Data Protection", second)
        self.assertEqual([("Data Protection", "zh-CN")], fake_client.calls)
        self.assertEqual(1, runtime.cache_hit_count)

    def test_runtime_returns_original_text_when_fail_fast_disabled(self):
        input_path = os.path.join(self.temp_dir.name, "sample.docx")
        with open(input_path, "wb") as handle:
            handle.write(b"doc")

        runtime = self.make_runtime(input_path, no_cache=True, resume=False, fail_fast=False)
        runtime.client = FailingClient("HTTP 429")

        result = runtime.translate_text("Inline CASB")

        self.assertEqual("Inline CASB", result)
        self.assertEqual(1, len(runtime.failed_texts))

    def test_state_manager_rejects_resume_when_input_changes(self):
        input_path = os.path.join(self.temp_dir.name, "sample.docx")
        with open(input_path, "wb") as handle:
            handle.write(b"v1")

        original_fingerprint = translator.fingerprint_file(input_path)
        manager = translator.TranslationStateManager(
            input_file=input_path,
            target_language="zh-CN",
            input_fingerprint=original_fingerprint,
            enabled=True,
        )
        manager.mark_completed("Data Protection")
        manager.save(status="running")

        with open(input_path, "wb") as handle:
            handle.write(b"v2")

        changed_fingerprint = translator.fingerprint_file(input_path)
        reloaded = translator.TranslationStateManager(
            input_file=input_path,
            target_language="zh-CN",
            input_fingerprint=changed_fingerprint,
            enabled=True,
        )

        self.assertFalse(reloaded.load())

    def test_ppt_paragraph_uses_paragraph_level_when_styles_match(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(0, 0, 100, 100)
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.text = ""
        paragraph.add_run().text = "Line speed "
        paragraph.add_run().text = "URL Filtering"

        self.assertFalse(translator.should_translate_ppt_paragraph_by_run(paragraph))

    def test_ppt_paragraph_uses_run_level_when_styles_differ(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(0, 0, 100, 100)
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.text = ""
        first = paragraph.add_run()
        first.text = "Inline "
        second = paragraph.add_run()
        second.text = "CASB"
        second.font.bold = True

        self.assertTrue(translator.should_translate_ppt_paragraph_by_run(paragraph))


if __name__ == "__main__":
    unittest.main()
