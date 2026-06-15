from __future__ import annotations

import argparse
import hashlib
import json
import os
import random
import re
import signal
import sys
import time
from dataclasses import dataclass, field
from typing import Callable

import docx
import requests
import urllib3
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

ENV_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
PROMPT_TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "llm_prompt.txt"
)
CACHE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "cache")
FONT_MODIFIED = "Microsoft YaHei"

DEFAULT_MODEL_NAME = "gemini-3-pro-high"
DEFAULT_ENDPOINT = "https://api.siliconflow.cn/v1"
DEFAULT_TEMPERATURE = 1.0
DEFAULT_ENABLE_THINKING = False
DEFAULT_REQUEST_TIMEOUT_SECONDS = 60.0
DEFAULT_MAX_RETRIES = 4
DEFAULT_INITIAL_BACKOFF_SECONDS = 1.0
DEFAULT_MAX_BACKOFF_SECONDS = 20.0
DEFAULT_MIN_REQUEST_INTERVAL_SECONDS = 0.75
DEFAULT_CONSECUTIVE_FAILURE_LIMIT = 3
DEFAULT_SAVE_PROGRESS_EVERY_N = 20
DEFAULT_RESUME_ENABLED = True

RUNTIME: "TranslationRuntime | None" = None


def load_env() -> dict[str, str]:
    config: dict[str, str] = {}
    if os.path.exists(ENV_PATH):
        with open(ENV_PATH, "r", encoding="utf-8") as handle:
            for raw_line in handle:
                line = raw_line.strip()
                if line and not line.startswith("#") and "=" in line:
                    key, value = line.split("=", 1)
                    config[key.strip()] = value.strip()
    return config


def load_prompt_template() -> str:
    if os.path.exists(PROMPT_TEMPLATE_PATH):
        with open(PROMPT_TEMPLATE_PATH, "r", encoding="utf-8") as handle:
            return handle.read()
    return "You are a professional translator. Translate to {target_language}."


PROMPT_TEMPLATE = load_prompt_template()


def get_prompt(target_language: str) -> str:
    return PROMPT_TEMPLATE.format(target_language=target_language)


def env_bool(env: dict[str, str], key: str, default: bool) -> bool:
    value = env.get(key)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def env_float(env: dict[str, str], key: str, default: float) -> float:
    value = env.get(key)
    if value is None:
        return default
    return float(value)


def env_int(env: dict[str, str], key: str, default: int) -> int:
    value = env.get(key)
    if value is None:
        return default
    return int(value)


def normalize_text_key(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def should_skip_translation(text: str) -> bool:
    normalized = normalize_text_key(text)
    if len(normalized) < 2:
        return True
    if re.fullmatch(r"[\W_]+", normalized):
        return True
    if re.fullmatch(r"(page|p\.?)?\s*\d+(\s*/\s*\d+)?", normalized, flags=re.I):
        return True
    if not re.search(r"[A-Za-z\u00C0-\u024F\u4e00-\u9fff]", normalized):
        return True
    return False


def fingerprint_file(path: str) -> str:
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


class TranslationAbort(RuntimeError):
    pass


class TranslationRequestError(RuntimeError):
    pass


@dataclass
class Config:
    api_key: str
    model_name: str
    endpoint: str
    temperature: float
    enable_thinking: bool
    request_timeout_seconds: float
    max_retries: int
    initial_backoff_seconds: float
    max_backoff_seconds: float
    min_request_interval_seconds: float
    consecutive_failure_limit: int
    save_progress_every_n: int
    resume_enabled: bool

    @property
    def is_groq_api(self) -> bool:
        return "groq.com" in self.endpoint.lower()

    @property
    def reasoning_effort(self) -> str:
        return "default" if self.enable_thinking else "none"


class LLMClient:
    def __init__(
        self,
        config: Config,
        fail_fast: bool,
        session: requests.Session | None = None,
        sleep_fn: Callable[[float], None] = time.sleep,
        monotonic_fn: Callable[[], float] = time.monotonic,
        random_fn: Callable[[], float] = random.random,
    ) -> None:
        self.config = config
        self.fail_fast = fail_fast
        self.session = session or requests.Session()
        self.sleep_fn = sleep_fn
        self.monotonic_fn = monotonic_fn
        self.random_fn = random_fn
        self.last_request_at: float | None = None
        self.api_request_count = 0
        self.retry_count = 0
        self.consecutive_failures = 0
        self.last_error = ""

    def translate(self, text: str, target_language: str) -> str:
        headers = {
            "Authorization": f"Bearer {self.config.api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": self.config.model_name,
            "messages": [
                {"role": "system", "content": get_prompt(target_language)},
                {"role": "user", "content": text},
            ],
            "temperature": self.config.temperature,
        }

        if self.config.is_groq_api:
            payload["reasoning_effort"] = self.config.reasoning_effort
        else:
            payload["enable_thinking"] = self.config.enable_thinking

        last_error = ""
        for attempt in range(self.config.max_retries + 1):
            self._wait_for_slot()
            self.api_request_count += 1
            try:
                response = self.session.post(
                    f"{self.config.endpoint}/chat/completions",
                    headers=headers,
                    json=payload,
                    timeout=(
                        self.config.request_timeout_seconds,
                        self.config.request_timeout_seconds,
                    ),
                    verify=False,
                )
                self.last_request_at = self.monotonic_fn()
                if response.status_code == 429 or 500 <= response.status_code < 600:
                    raise TranslationRequestError(
                        self._http_error_message(response.status_code, response.text)
                    )
                response.raise_for_status()
                translated_text = (
                    response.json()["choices"][0]["message"]["content"].strip()
                )
                self.consecutive_failures = 0
                self.last_error = ""
                return translated_text
            except requests.exceptions.HTTPError as exc:
                last_error = str(exc)
                if attempt >= self.config.max_retries:
                    break
                self.retry_count += 1
                delay = self._retry_delay(attempt, None)
                print(
                    f"HTTP error, retrying in {delay:.1f}s "
                    f"({attempt + 1}/{self.config.max_retries}): {last_error}"
                )
                self.sleep_fn(delay)
            except TranslationRequestError as exc:
                last_error = str(exc)
                if attempt >= self.config.max_retries:
                    break
                self.retry_count += 1
                delay = self._retry_delay(attempt, getattr(locals().get("response"), "headers", None))
                print(
                    f"Retryable API error, retrying in {delay:.1f}s "
                    f"({attempt + 1}/{self.config.max_retries}): {last_error}"
                )
                self.sleep_fn(delay)
            except (
                requests.exceptions.ConnectionError,
                requests.exceptions.SSLError,
                requests.exceptions.Timeout,
            ) as exc:
                last_error = str(exc)
                if attempt >= self.config.max_retries:
                    break
                self.retry_count += 1
                delay = self._retry_delay(attempt, None)
                print(
                    f"Network error, retrying in {delay:.1f}s "
                    f"({attempt + 1}/{self.config.max_retries}): {last_error}"
                )
                self.sleep_fn(delay)
            except (KeyError, ValueError, json.JSONDecodeError) as exc:
                last_error = f"Invalid API response: {exc}"
                break

        self.consecutive_failures += 1
        self.last_error = last_error or "Unknown translation failure"
        if self.fail_fast or (
            self.consecutive_failures >= self.config.consecutive_failure_limit
        ):
            raise TranslationAbort(self.last_error)
        raise TranslationRequestError(self.last_error)

    def _wait_for_slot(self) -> None:
        if self.last_request_at is None:
            return
        elapsed = self.monotonic_fn() - self.last_request_at
        if elapsed >= self.config.min_request_interval_seconds:
            return
        self.sleep_fn(self.config.min_request_interval_seconds - elapsed)

    def _retry_delay(self, attempt: int, headers: dict[str, str] | None) -> float:
        retry_after = parse_retry_after(headers)
        if retry_after is not None:
            return retry_after
        delay = min(
            self.config.max_backoff_seconds,
            self.config.initial_backoff_seconds * (2**attempt),
        )
        jitter = delay * 0.1 * self.random_fn()
        return delay + jitter

    @staticmethod
    def _http_error_message(status_code: int, body: str) -> str:
        body_preview = normalize_text_key(body)[:200]
        if body_preview:
            return f"HTTP {status_code}: {body_preview}"
        return f"HTTP {status_code}"


def parse_retry_after(headers: dict[str, str] | None) -> float | None:
    if not headers:
        return None
    raw_value = headers.get("Retry-After")
    if not raw_value:
        return None
    try:
        return max(0.0, float(raw_value))
    except ValueError:
        return None


@dataclass
class TranslationStateManager:
    input_file: str
    target_language: str
    input_fingerprint: str
    enabled: bool
    unique_text_total: int = 0
    completed_keys: set[str] = field(default_factory=set)
    last_completed_key: str = ""
    status: str = "running"

    def __post_init__(self) -> None:
        os.makedirs(CACHE_DIR, exist_ok=True)
        state_seed = (
            f"{os.path.abspath(self.input_file)}|{self.target_language}".encode("utf-8")
        )
        state_id = hashlib.sha1(state_seed).hexdigest()[:12]
        self.state_file = os.path.join(CACHE_DIR, f"{state_id}.translation-state.json")

    def load(self) -> bool:
        if not self.enabled or not os.path.exists(self.state_file):
            return False
        with open(self.state_file, "r", encoding="utf-8") as handle:
            data = json.load(handle)
        if data.get("input_file") != os.path.abspath(self.input_file):
            return False
        if data.get("target_language") != self.target_language:
            return False
        if data.get("input_fingerprint") != self.input_fingerprint:
            print("Existing resume state ignored because the input file changed.")
            return False
        self.completed_keys = set(data.get("completed_keys", []))
        self.last_completed_key = data.get("last_completed_key", "")
        self.unique_text_total = data.get("unique_text_total", self.unique_text_total)
        self.status = data.get("status", "running")
        return True

    def mark_completed(self, cache_key: str) -> None:
        self.completed_keys.add(cache_key)
        self.last_completed_key = cache_key

    def save(self, status: str) -> None:
        if not self.enabled:
            return
        payload = {
            "input_file": os.path.abspath(self.input_file),
            "target_language": self.target_language,
            "input_fingerprint": self.input_fingerprint,
            "unique_text_total": self.unique_text_total,
            "completed_keys": sorted(self.completed_keys),
            "last_completed_key": self.last_completed_key,
            "translated_unique_count": len(self.completed_keys),
            "status": status,
            "saved_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        temp_path = f"{self.state_file}.tmp"
        with open(temp_path, "w", encoding="utf-8") as handle:
            json.dump(payload, handle, ensure_ascii=False, indent=2)
        os.replace(temp_path, self.state_file)
        self.status = status

    def clear(self) -> None:
        if self.enabled and os.path.exists(self.state_file):
            os.remove(self.state_file)


class TranslationRuntime:
    def __init__(
        self,
        args: argparse.Namespace,
        config: Config,
        input_file: str,
        target_language: str,
    ) -> None:
        self.args = args
        self.config = config
        self.input_file = input_file
        self.target_language = target_language
        self.translation_cache, self.cache_file = init_cache(
            target_language=target_language,
            no_cache=args.no_cache,
        )
        self.client = LLMClient(config=config, fail_fast=args.fail_fast)
        self.input_fingerprint = fingerprint_file(input_file)
        self.state = TranslationStateManager(
            input_file=input_file,
            target_language=target_language,
            input_fingerprint=self.input_fingerprint,
            enabled=args.resume,
        )
        self.cache_hit_count = 0
        self.new_translation_count = 0
        self.skipped_count = 0
        self.failed_texts: list[dict[str, str]] = []
        self.unique_candidate_keys: set[str] = set()
        self._pending_checkpoint_count = 0

    def collect_candidate(self, text: str) -> None:
        if should_skip_translation(text):
            return
        self.unique_candidate_keys.add(normalize_text_key(text))

    def load_resume_state(self) -> None:
        self.state.unique_text_total = len(self.unique_candidate_keys)
        if self.state.load():
            print(
                "Resume state loaded: "
                f"{len(self.state.completed_keys)}/{self.state.unique_text_total} unique texts completed."
            )
        else:
            self.state.unique_text_total = len(self.unique_candidate_keys)
        print(f"Unique translatable texts this run: {len(self.unique_candidate_keys)}")

    def translate_text(self, text: str) -> str:
        if should_skip_translation(text):
            self.skipped_count += 1
            return text

        cache_key = normalize_text_key(text)
        if cache_key in self.translation_cache:
            self.cache_hit_count += 1
            self.state.mark_completed(cache_key)
            print(f"[Cache hit] {cache_key[:40]}{'...' if len(cache_key) > 40 else ''}")
            return self.translation_cache[cache_key]

        try:
            translated_text = self.client.translate(text, self.target_language)
        except TranslationRequestError as exc:
            self.record_failure(text, str(exc))
            print(f"Translation error: {exc}")
            return text
        self.translation_cache[cache_key] = translated_text
        self.state.mark_completed(cache_key)
        self.new_translation_count += 1
        self._pending_checkpoint_count += 1
        self.maybe_checkpoint()
        return translated_text

    def maybe_checkpoint(self, force: bool = False, status: str = "running") -> None:
        if not force and self._pending_checkpoint_count < self.config.save_progress_every_n:
            return
        save_cache(self.translation_cache, self.cache_file, self.args.no_cache)
        self.state.save(status=status)
        self._pending_checkpoint_count = 0

    def record_failure(self, text: str, error: str) -> None:
        self.failed_texts.append({"text": normalize_text_key(text), "error": error})

    def finalize(self, success: bool) -> None:
        save_cache(self.translation_cache, self.cache_file, self.args.no_cache)
        if success:
            self.state.clear()
        else:
            self.state.save(status="aborted")

    def print_summary(self) -> None:
        print(f"Unique texts: {len(self.unique_candidate_keys)}")
        print(f"Cache hits this run: {self.cache_hit_count}")
        print(f"New translations this run: {self.new_translation_count}")
        print(f"Skipped texts this run: {self.skipped_count}")
        print(f"Actual API requests: {self.client.api_request_count}")
        print(f"Retry count: {self.client.retry_count}")
        if self.client.last_error:
            print(f"Last error: {self.client.last_error}")
        if self.failed_texts:
            print(f"Failed texts: {len(self.failed_texts)}")


def init_cache(target_language: str, no_cache: bool) -> tuple[dict[str, str], str | None]:
    os.makedirs(CACHE_DIR, exist_ok=True)
    cache_file = os.path.join(CACHE_DIR, f"global-{target_language}.json")
    if no_cache or not os.path.exists(cache_file):
        return {}, cache_file
    try:
        with open(cache_file, "r", encoding="utf-8") as handle:
            cache = json.load(handle)
        print(
            f"Loaded {len(cache)} translation cache entries "
            f"(global, target language: {target_language})"
        )
        return cache, cache_file
    except Exception as exc:
        print(f"Failed to load cache: {exc}")
        return {}, cache_file


def save_cache(
    translation_cache: dict[str, str], cache_file: str | None, no_cache: bool
) -> None:
    if no_cache or not cache_file:
        return
    try:
        with open(cache_file, "w", encoding="utf-8") as handle:
            json.dump(translation_cache, handle, ensure_ascii=False, indent=2)
        print(
            f"Saved {len(translation_cache)} translation cache entries "
            f"(global, target language: {os.path.basename(cache_file)})"
        )
    except Exception as exc:
        print(f"Failed to save cache: {exc}")


def safe_set_font_color(run) -> None:
    try:
        pass
    except AttributeError:
        pass


def safe_set_font(run) -> None:
    try:
        run.font.name = FONT_MODIFIED
    except AttributeError:
        pass


def build_parser(env_config: dict[str, str]) -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Translate PowerPoint or Word files")
    parser.add_argument("input_file", nargs="?", help="Input PPT or Word file")
    parser.add_argument(
        "target_language",
        nargs="?",
        default="zh-CN",
        help="Target language code (default: zh-CN)",
    )
    parser.add_argument(
        "--type", choices=["ppt", "word"], help="Specify file type (ppt or word)"
    )
    parser.add_argument("--no-cache", action="store_true", help="Disable translation cache")
    parser.add_argument(
        "--resume",
        dest="resume",
        action="store_true",
        help="Resume from saved progress state when available",
    )
    parser.add_argument(
        "--no-resume",
        dest="resume",
        action="store_false",
        help="Disable resume state loading and saving",
    )
    parser.add_argument(
        "--fail-fast",
        dest="fail_fast",
        action="store_true",
        help="Abort after a translation unit exhausts its retries",
    )
    parser.add_argument(
        "--no-fail-fast",
        dest="fail_fast",
        action="store_false",
        help="Keep the original text after exhausted retries and continue",
    )
    parser.add_argument(
        "--min-request-interval",
        type=float,
        default=None,
        help="Minimum interval in seconds between API requests",
    )
    parser.add_argument(
        "--max-retries",
        type=int,
        default=None,
        help="Maximum retry count for retryable API errors",
    )
    parser.set_defaults(
        resume=env_bool(env_config, "RESUME_ENABLED", DEFAULT_RESUME_ENABLED),
        fail_fast=True,
    )
    return parser


def build_config(env_config: dict[str, str], args: argparse.Namespace) -> Config:
    return Config(
        api_key=env_config.get("LLM_API_KEY", ""),
        model_name=env_config.get("MODEL_NAME", DEFAULT_MODEL_NAME),
        endpoint=env_config.get("ENDPOINT", DEFAULT_ENDPOINT),
        temperature=env_float(env_config, "TEMPERATURE", DEFAULT_TEMPERATURE),
        enable_thinking=env_bool(
            env_config, "ENABLE_THINKING", DEFAULT_ENABLE_THINKING
        ),
        request_timeout_seconds=env_float(
            env_config, "REQUEST_TIMEOUT_SECONDS", DEFAULT_REQUEST_TIMEOUT_SECONDS
        ),
        max_retries=(
            args.max_retries
            if args.max_retries is not None
            else env_int(env_config, "MAX_RETRIES", DEFAULT_MAX_RETRIES)
        ),
        initial_backoff_seconds=env_float(
            env_config,
            "INITIAL_BACKOFF_SECONDS",
            DEFAULT_INITIAL_BACKOFF_SECONDS,
        ),
        max_backoff_seconds=env_float(
            env_config,
            "MAX_BACKOFF_SECONDS",
            DEFAULT_MAX_BACKOFF_SECONDS,
        ),
        min_request_interval_seconds=(
            args.min_request_interval
            if args.min_request_interval is not None
            else env_float(
                env_config,
                "MIN_REQUEST_INTERVAL_SECONDS",
                DEFAULT_MIN_REQUEST_INTERVAL_SECONDS,
            )
        ),
        consecutive_failure_limit=env_int(
            env_config,
            "CONSECUTIVE_FAILURE_LIMIT",
            DEFAULT_CONSECUTIVE_FAILURE_LIMIT,
        ),
        save_progress_every_n=env_int(
            env_config,
            "SAVE_PROGRESS_EVERY_N",
            DEFAULT_SAVE_PROGRESS_EVERY_N,
        ),
        resume_enabled=env_bool(env_config, "RESUME_ENABLED", DEFAULT_RESUME_ENABLED),
    )


def resolve_input_file(args: argparse.Namespace) -> tuple[str, str]:
    if args.target_language and (
        args.target_language.startswith(".\\") or args.target_language.startswith("./")
    ):
        input_file = args.target_language
        args.target_language = "zh-CN"
        args.input_file = input_file
        print(
            f"Warning: Argument '{input_file}' looks like a file path, "
            "not a language code. Using it as the input file."
        )

    if args.input_file:
        input_file = args.input_file
        print(f"Using specified input file: {input_file}")
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"File not found: {input_file}")
        file_ext = os.path.splitext(input_file)[1].lower()
        if args.type:
            file_type = args.type
        elif file_ext in [".pptx", ".ppt"]:
            file_type = "ppt"
        elif file_ext in [".docx", ".doc"]:
            file_type = "word"
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        return input_file, file_type

    if args.type == "word":
        doc_files = [f for f in os.listdir(".") if f.endswith((".docx", ".doc"))]
        if not doc_files:
            raise FileNotFoundError("No .docx or .doc files found in current directory.")
        return doc_files[0], "word"

    ppt_files = [f for f in os.listdir(".") if f.endswith((".pptx", ".ppt"))]
    if not ppt_files:
        raise FileNotFoundError("No .pptx or .ppt files found in current directory.")
    return ppt_files[0], "ppt"


def output_path_for(input_file: str, target_language: str) -> str:
    root, ext = os.path.splitext(input_file)
    return f"{root}-{target_language}{ext}"


def run_style_signature(run) -> tuple:
    font = getattr(run, "font", None)
    if font is None:
        return ("", None, None, None, None, "")
    color_rgb = ""
    try:
        if font.color is not None and font.color.rgb is not None:
            color_rgb = str(font.color.rgb)
    except Exception:
        color_rgb = ""
    return (
        font.name,
        font.size,
        font.bold,
        font.italic,
        font.underline,
        color_rgb,
    )


def should_translate_ppt_paragraph_by_run(paragraph) -> bool:
    non_empty_runs = [run for run in paragraph.runs if run.text and run.text.strip()]
    if len(non_empty_runs) <= 1:
        return False
    if any(re.search(r"(\{\{.*?\}\}|\{.*?\}|<[^>]+>)", run.text) for run in non_empty_runs):
        return True
    signatures = {run_style_signature(run) for run in non_empty_runs}
    return len(signatures) > 1


def set_ppt_paragraph_text(paragraph, translated_text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = translated_text
        safe_set_font(paragraph.runs[0])
        safe_set_font_color(paragraph.runs[0])
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = translated_text


def collect_ppt_text_frame(text_frame, runtime: TranslationRuntime) -> None:
    for paragraph in text_frame.paragraphs:
        if should_translate_ppt_paragraph_by_run(paragraph):
            for run in paragraph.runs:
                runtime.collect_candidate(run.text)
        else:
            runtime.collect_candidate(paragraph.text)


def translate_ppt_text_frame(text_frame, runtime: TranslationRuntime) -> None:
    for paragraph in text_frame.paragraphs:
        if should_translate_ppt_paragraph_by_run(paragraph):
            for run in paragraph.runs:
                original_text = run.text
                if not original_text:
                    continue
                print(f"Original text: {original_text}")
                translated_text = runtime.translate_text(original_text)
                run.text = translated_text
                safe_set_font(run)
                safe_set_font_color(run)
                print(f"Updated text: {run.text}")
            continue

        original_text = paragraph.text
        if not original_text:
            continue
        print(f"Original text: {original_text}")
        translated_text = runtime.translate_text(original_text)
        set_ppt_paragraph_text(paragraph, translated_text)
        print(f"Updated text: {translated_text}")


def collect_table(table, runtime: TranslationRuntime) -> None:
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                collect_ppt_text_frame(cell.text_frame, runtime)


def translate_table(table, runtime: TranslationRuntime) -> None:
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                translate_ppt_text_frame(cell.text_frame, runtime)


def collect_chart(chart, runtime: TranslationRuntime) -> None:
    if chart.has_title:
        runtime.collect_candidate(chart.chart_title.text_frame.text)
    try:
        if chart.chart_data and isinstance(chart.chart_data, CategoryChartData):
            for category in chart.chart_data.categories:
                runtime.collect_candidate(category.label)
            for series in chart.chart_data.series:
                runtime.collect_candidate(series.name)
    except Exception:
        pass


def translate_chart(chart, runtime: TranslationRuntime) -> None:
    if chart.has_title:
        chart.chart_title.text_frame.text = runtime.translate_text(
            chart.chart_title.text_frame.text
        )
    try:
        if chart.chart_data and isinstance(chart.chart_data, CategoryChartData):
            for category in chart.chart_data.categories:
                category.label = runtime.translate_text(category.label)
            for series in chart.chart_data.series:
                series.name = runtime.translate_text(series.name)
    except Exception:
        pass


def iter_complex_shape_text_elements(shape):
    for element in shape.element.iter():
        if element.tag.endswith("}t") and element.text and element.text.strip():
            yield element


def collect_shape(shape, runtime: TranslationRuntime) -> None:
    if shape.has_text_frame:
        collect_ppt_text_frame(shape.text_frame, runtime)
    elif shape.has_table:
        collect_table(shape.table, runtime)
    elif shape.has_chart:
        collect_chart(shape.chart, runtime)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            collect_shape(child_shape, runtime)
    elif hasattr(shape, "graphic") and hasattr(shape.graphic, "graphicData"):
        for element in iter_complex_shape_text_elements(shape):
            runtime.collect_candidate(element.text)


def translate_shape(shape, runtime: TranslationRuntime) -> None:
    if shape.has_text_frame:
        translate_ppt_text_frame(shape.text_frame, runtime)
    elif shape.has_table:
        translate_table(shape.table, runtime)
    elif shape.has_chart:
        translate_chart(shape.chart, runtime)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            translate_shape(child_shape, runtime)
    elif hasattr(shape, "graphic") and hasattr(shape.graphic, "graphicData"):
        for element in iter_complex_shape_text_elements(shape):
            element.text = runtime.translate_text(element.text)


def collect_ppt_texts(prs: Presentation, runtime: TranslationRuntime) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            collect_shape(shape, runtime)
        if hasattr(slide, "header"):
            collect_ppt_text_frame(slide.header.text_frame, runtime)
        if hasattr(slide, "footer"):
            collect_ppt_text_frame(slide.footer.text_frame, runtime)
    if prs.core_properties.title:
        runtime.collect_candidate(prs.core_properties.title)
    if prs.core_properties.subject:
        runtime.collect_candidate(prs.core_properties.subject)


def translate_pptx(input_file: str, output_file: str, runtime: TranslationRuntime) -> None:
    prs = Presentation(input_file)
    collect_ppt_texts(prs, runtime)
    runtime.load_resume_state()

    for slide in prs.slides:
        for shape in slide.shapes:
            translate_shape(shape, runtime)
        if hasattr(slide, "header"):
            translate_ppt_text_frame(slide.header.text_frame, runtime)
        if hasattr(slide, "footer"):
            translate_ppt_text_frame(slide.footer.text_frame, runtime)

    if prs.core_properties.title:
        prs.core_properties.title = runtime.translate_text(prs.core_properties.title)
    if prs.core_properties.subject:
        prs.core_properties.subject = runtime.translate_text(prs.core_properties.subject)

    prs.save(output_file)
    print(f"Translated PPT saved to {output_file}")


def collect_word_paragraph(paragraph, runtime: TranslationRuntime) -> None:
    runtime.collect_candidate(paragraph.text)


def translate_word_paragraph(paragraph, runtime: TranslationRuntime) -> None:
    full_text = paragraph.text
    if not full_text.strip():
        return

    translated_text = runtime.translate_text(full_text)
    for _ in range(len(paragraph.runs)):
        paragraph._element.remove(paragraph._element.r_lst[0])
    new_run = paragraph.add_run(translated_text)
    safe_set_font(new_run)
    print(f"Original: {full_text[:50]}...")
    print(f"Translated: {translated_text[:50]}...")


def collect_word_table(table, runtime: TranslationRuntime) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                collect_word_paragraph(paragraph, runtime)


def translate_word_table(table, runtime: TranslationRuntime) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                translate_word_paragraph(paragraph, runtime)


def collect_docx_texts(doc: docx.Document, runtime: TranslationRuntime) -> None:
    for paragraph in doc.paragraphs:
        collect_word_paragraph(paragraph, runtime)
    for table in doc.tables:
        collect_word_table(table, runtime)
    for section in doc.sections:
        for header in section.header.paragraphs:
            collect_word_paragraph(header, runtime)
        for footer in section.footer.paragraphs:
            collect_word_paragraph(footer, runtime)
    if hasattr(doc.core_properties, "title") and doc.core_properties.title:
        runtime.collect_candidate(doc.core_properties.title)
    if hasattr(doc.core_properties, "subject") and doc.core_properties.subject:
        runtime.collect_candidate(doc.core_properties.subject)


def translate_docx(input_file: str, output_file: str, runtime: TranslationRuntime) -> None:
    print(f"Translating Word document: {input_file}")
    doc = docx.Document(input_file)
    collect_docx_texts(doc, runtime)
    runtime.load_resume_state()

    for paragraph in doc.paragraphs:
        translate_word_paragraph(paragraph, runtime)
    for table in doc.tables:
        translate_word_table(table, runtime)
    for section in doc.sections:
        for header in section.header.paragraphs:
            translate_word_paragraph(header, runtime)
        for footer in section.footer.paragraphs:
            translate_word_paragraph(footer, runtime)
    if hasattr(doc.core_properties, "title") and doc.core_properties.title:
        doc.core_properties.title = runtime.translate_text(doc.core_properties.title)
    if hasattr(doc.core_properties, "subject") and doc.core_properties.subject:
        doc.core_properties.subject = runtime.translate_text(doc.core_properties.subject)

    doc.save(output_file)
    print(f"Translated Word document saved to {output_file}")


def install_signal_handler(runtime: TranslationRuntime) -> None:
    def signal_handler(sig, frame) -> None:
        runtime.finalize(success=False)
        runtime.print_summary()
        raise SystemExit(130)

    signal.signal(signal.SIGINT, signal_handler)


def validate_config(config: Config) -> None:
    if not config.api_key:
        raise ValueError("LLM_API_KEY is required. Please configure it in .env.")


def main(argv: list[str] | None = None) -> int:
    global RUNTIME

    env_config = load_env()
    parser = build_parser(env_config)
    args = parser.parse_args(argv)
    config = build_config(env_config, args)
    validate_config(config)

    input_file, file_type = resolve_input_file(args)
    output_file = output_path_for(input_file, args.target_language)
    runtime = TranslationRuntime(
        args=args,
        config=config,
        input_file=input_file,
        target_language=args.target_language,
    )
    RUNTIME = runtime
    install_signal_handler(runtime)

    try:
        if file_type == "ppt":
            print(f"Translating PPT file '{input_file}' to {args.target_language}")
            translate_pptx(input_file=input_file, output_file=output_file, runtime=runtime)
        else:
            print(f"Translating Word file '{input_file}' to {args.target_language}")
            translate_docx(input_file=input_file, output_file=output_file, runtime=runtime)
        runtime.finalize(success=True)
        runtime.print_summary()
        return 0
    except TranslationAbort as exc:
        runtime.record_failure("<aborted>", str(exc))
        runtime.finalize(success=False)
        runtime.print_summary()
        print(f"Translation aborted: {exc}")
        return 1
    except Exception as exc:
        runtime.record_failure("<unexpected>", str(exc))
        runtime.finalize(success=False)
        runtime.print_summary()
        print(f"Translation failed: {exc}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
